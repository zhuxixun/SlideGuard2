from __future__ import annotations

import csv
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from slideguard.preview.text_layout import measure_single_line


@dataclass(frozen=True, slots=True)
class MetricCase:
    case_id: str
    text: str
    font_size_pt: float
    bold: bool
    width_delta_pt: float


CASES = (
    MetricCase("cn_exact_minus_3", "微软雅黑文本溢出边界", 24, False, -3),
    MetricCase("cn_exact_minus_2", "微软雅黑文本溢出边界", 24, False, -2),
    MetricCase("cn_exact_minus_1", "微软雅黑文本溢出边界", 24, False, -1),
    MetricCase("cn_exact", "微软雅黑文本溢出边界", 24, False, 0),
    MetricCase("cn_exact_plus_2", "微软雅黑文本溢出边界", 24, False, 2),
    MetricCase("mixed_exact", "SlideGuard 文本 Overflow", 24, False, 0),
    MetricCase("bold_exact", "加粗标题：文本溢出", 24, True, 0),
    MetricCase("body_14pt", "正文最小字号十四点", 14, False, 0),
)


def generate_text_metric_probe(output_dir: Path) -> tuple[Path, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / "text-metric-probe.pptx"
    csv_path = output_dir / "text-metric-probe.csv"
    if pptx_path.exists() or csv_path.exists():
        raise FileExistsError("测量探针输出已存在，请先移动或删除旧文件")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    rows: list[dict[str, str | float]] = []

    for case in CASES:
        measurement = measure_single_line(
            case.text,
            font_size_pt=case.font_size_pt,
            bold=case.bold,
        )
        target_width = max(1, measurement.width_pt + case.width_delta_pt)
        slide = presentation.slides.add_slide(blank_layout)
        _add_label(slide, case, measurement.width_pt, target_width)
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(2),
            Pt(target_width),
            Pt(measurement.height_pt + 4),
        )
        frame = box.text_frame
        frame.clear()
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        frame.word_wrap = False
        frame.auto_size = MSO_AUTO_SIZE.NONE
        run = frame.paragraphs[0].add_run()
        run.text = case.text
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(case.font_size_pt)
        run.font.bold = case.bold
        rows.append(
            {
                "case_id": case.case_id,
                "slide": len(presentation.slides),
                "text": case.text,
                "font_size_pt": case.font_size_pt,
                "bold": str(case.bold).lower(),
                "pillow_width_pt": round(measurement.width_pt, 3),
                "pillow_height_pt": round(measurement.height_pt, 3),
                "box_width_pt": round(target_width, 3),
                "width_delta_pt": case.width_delta_pt,
                "wps_bound_width_pt": "",
                "wps_absolute_error_pt": "",
                "wps_result": "",
                "powerpoint_visible_width_pt": "",
                "absolute_error_pt": "",
                "review_result": "",
            }
        )

    presentation.save(pptx_path)
    with csv_path.open("w", encoding="utf-8-sig", newline="") as stream:
        writer = csv.DictWriter(stream, fieldnames=list(rows[0]))
        writer.writeheader()
        writer.writerows(rows)
    return pptx_path, csv_path


def _add_label(slide, case: MetricCase, measured: float, target: float) -> None:  # type: ignore[no-untyped-def]
    label = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    paragraph = label.text_frame.paragraphs[0]
    paragraph.text = (
        f"{case.case_id} | Pillow={measured:.3f}pt | "
        f"文本框={target:.3f}pt | 差值={case.width_delta_pt:+.1f}pt"
    )
    paragraph.runs[0].font.name = "Microsoft YaHei"
    paragraph.runs[0].font.size = Pt(18)
    paragraph.runs[0].font.color.rgb = RGBColor(192, 0, 0)
