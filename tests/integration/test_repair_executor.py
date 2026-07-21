from dataclasses import replace
from pathlib import Path
import shutil
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from slideguard.pptx.importer import inspect_pptx
from slideguard.repair.executor import execute_and_recheck, execute_fix_plan
from slideguard.repair.planner import build_fix_plan
from slideguard.rules.factory import issue
from slideguard.rules.models import IssueStatus, Severity
from slideguard.scan.models import ScanMode, ScanRequest
from slideguard.scan.orchestrator import run_scan


def _standard_result(path: Path, issues):  # type: ignore[no-untyped-def]
    result = run_scan(
        inspect_pptx(path),
        ScanRequest(ScanMode.CUSTOM, selected_rules=("R002",)),
        rules={"R002": lambda *_: tuple(issues)},
    )
    return replace(result, mode=ScanMode.STANDARD)


def _fixable(rule_id: str, fact: str, key: str, kind: str, target: str):  # type: ignore[no-untyped-def]
    return issue(
        fact_key=fact,
        rule_id=rule_id,
        slide_index=1,
        object_keys=(key,),
        severity=Severity.S3,
        actual_value="old",
        expected_value=target,
        evidence="evidence",
        suggestion="suggestion",
        can_auto_fix=True,
        fix_kind=kind,
        fix_target=target,
    )


def test_executor_applies_font_size_and_position_with_minimal_part_changes(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "fixed.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Test"
    run.font.name = "Arial"
    run.font.size = Pt(12)
    key = f"ppt/slides/slide1.xml:shape:{shape.shape_id}"
    document.save(source)
    issues = (
        _fixable("R004", f"R004:1:{key}:font:0:4:Arial", key, "replace_font", "Microsoft YaHei"),
        _fixable("R005", f"R005:1:{key}:minimum-size:0:4", key, "replace_font_size", "14"),
        _fixable("R007", f"R007:1:{key}:alignment:left", key, "move_x", "144"),
    )
    result = _standard_result(source, issues)
    plan = build_fix_plan(result, tuple(found.issue_id for found in issues), output)
    execute_fix_plan(plan)

    reopened = Presentation(output)
    fixed = reopened.slides[0].shapes[0]
    fixed_run = fixed.text_frame.paragraphs[0].runs[0]
    assert fixed.left == Inches(2)
    assert fixed_run.font.name == "Microsoft YaHei"
    assert fixed_run.font.size == Pt(14)
    with ZipFile(source) as before, ZipFile(output) as after:
        unchanged = set(before.namelist()) - {"ppt/slides/slide1.xml"}
        assert all(before.read(name) == after.read(name) for name in unchanged)


def test_executor_sets_title_style_and_splits_suffix_then_rechecks(tmp_path: Path) -> None:
    source = tmp_path / "title.pptx"
    output = tmp_path / "title-fixed.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[0])
    title = slide.shapes.title
    title.text = "主题：说明文字"
    key = f"ppt/slides/slide1.xml:shape:{title.shape_id}"
    document.save(source)
    issues = (
        _fixable("R009", f"selected-title-style:{key}", key, "set_title_style", "style"),
        _fixable("R009", f"selected-title-overflow:{key}", key, "scale_title_suffix", "18"),
    )
    result = _standard_result(source, issues)
    plan = build_fix_plan(result, tuple(found.issue_id for found in issues), output)
    repaired = execute_and_recheck(plan)

    reopened = Presentation(output)
    runs = reopened.slides[0].shapes.title.text_frame.paragraphs[0].runs
    assert [run.text for run in runs] == ["主题：", "说明文字"]
    assert runs[0].font.name == "Microsoft YaHei"
    assert runs[0].font.size == Pt(24)
    assert runs[0].font.bold is True
    assert runs[0].font.color.rgb == RGBColor(192, 0, 0)
    assert runs[1].font.size == Pt(18)
    assert repaired.destination == output.resolve()
    assert set(repaired.fixed_issue_ids) == {found.issue_id for found in issues}


def test_executor_marks_selected_issue_as_fix_failed_when_recheck_still_finds_it(
    tmp_path: Path, monkeypatch
) -> None:  # type: ignore[no-untyped-def]
    source = tmp_path / "source.pptx"
    output = tmp_path / "unchanged.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    shape.text = "Test"
    document.save(source)
    key = f"ppt/slides/slide1.xml:shape:{shape.shape_id}"
    found = _fixable("R004", f"R004:1:{key}:font:0:4:Arial", key, "replace_font", "Microsoft YaHei")
    baseline = _standard_result(source, (found,))
    plan = build_fix_plan(baseline, (found.issue_id,), output)

    monkeypatch.setattr(
        "slideguard.repair.executor.execute_fix_plan",
        lambda _: shutil.copyfile(source, output),
    )
    monkeypatch.setattr("slideguard.repair.executor.run_scan", lambda *_args, **_kwargs: baseline)

    repaired = execute_and_recheck(plan)

    assert repaired.unresolved_issue_ids == (found.issue_id,)
    assert repaired.fixed_issue_ids == ()
    assert repaired.verification_scan.issues[0].status is IssueStatus.FIX_FAILED
