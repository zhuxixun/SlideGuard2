from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.util import Inches, Pt

from slideguard.pptx.patcher import XmlAttributePatch, patch_pptx


def test_real_pptx_font_and_position_patch_reopens_with_python_pptx(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "output.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "需要修复的标题"
    run.font.name = "Arial"
    run.font.size = Pt(24)
    shape_id = shape.shape_id
    presentation.save(source)

    shape_xpath = f"//p:sp[p:nvSpPr/p:cNvPr[@id='{shape_id}']]"
    patch_pptx(
        source,
        output,
        (
            XmlAttributePatch(
                "ppt/slides/slide1.xml",
                f"{shape_xpath}/p:spPr/a:xfrm/a:off",
                (("x", str(Inches(2))), ("y", str(Inches(1.5)))),
            ),
            XmlAttributePatch(
                "ppt/slides/slide1.xml",
                f"{shape_xpath}/p:txBody/a:p/a:r/a:rPr/a:latin",
                (("typeface", "Microsoft YaHei"),),
            ),
        ),
    )

    reopened = Presentation(output)
    reopened_shape = reopened.slides[0].shapes[0]
    assert reopened_shape.left == Inches(2)
    assert reopened_shape.top == Inches(1.5)
    assert reopened_shape.text_frame.paragraphs[0].runs[0].font.name == "Microsoft YaHei"

    with ZipFile(source) as before, ZipFile(output) as after:
        unchanged = set(before.namelist()) - {"ppt/slides/slide1.xml"}
        assert all(before.read(name) == after.read(name) for name in unchanged)

