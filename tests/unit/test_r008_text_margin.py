from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx.snapshot import build_snapshot
from slideguard.rules.models import Severity
from slideguard.rules.r008_text_margin import check_text_margins


def _add_text(slide, left, top, margin_left=0):  # type: ignore[no-untyped-def]
    box = slide.shapes.add_textbox(left, top, Inches(3), Inches(1))
    box.text_frame.margin_left = margin_left
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "安全边距"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(18)
    return box


def test_r008_uses_visible_text_not_shape_outer_bounds(tmp_path: Path) -> None:
    path = tmp_path / "actual-text-bounds.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    _add_text(slide, 0, Inches(1), margin_left=Pt(30))
    document.save(path)
    assert check_text_margins(build_snapshot(inspect_pptx(path))) == ()


def test_r008_reports_text_inside_three_percent_margin(tmp_path: Path) -> None:
    path = tmp_path / "unsafe.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    _add_text(slide, 0, Inches(1))
    document.save(path)
    issues = check_text_margins(build_snapshot(inspect_pptx(path)))
    assert len(issues) == 1
    assert issues[0].severity is Severity.S3
    assert "左侧" in issues[0].evidence
    assert issues[0].can_auto_fix is False


def test_r008_accepts_exact_three_percent_boundary(tmp_path: Path) -> None:
    path = tmp_path / "boundary.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    safe_left = document.slide_width * 3 // 100
    safe_top = document.slide_height * 3 // 100
    _add_text(slide, safe_left, safe_top)
    document.save(path)
    assert check_text_margins(build_snapshot(inspect_pptx(path))) == ()


def test_r008_checks_table_cell_text(tmp_path: Path) -> None:
    path = tmp_path / "table-margin.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    cell = slide.shapes.add_table(1, 1, 0, Inches(1), Inches(2), Inches(1)).table.cell(0, 0)
    cell.margin_left = 0
    run = cell.text_frame.paragraphs[0].add_run()
    run.text = "表格文字"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(18)
    document.save(path)
    issues = check_text_margins(build_snapshot(inspect_pptx(path)))
    assert len(issues) == 1
    assert ":cell:0:0" in issues[0].object_keys[0]
