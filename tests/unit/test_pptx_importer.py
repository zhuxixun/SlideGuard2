from pathlib import Path
from zipfile import ZipFile

import pytest
from pptx import Presentation

from slideguard.pptx.importer import PptxImportError, inspect_pptx


def test_inspect_pptx_returns_immutable_metadata(tmp_path: Path) -> None:
    path = tmp_path / "sample.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.slides.add_slide(document.slide_layouts[6])
    document.save(path)
    result = inspect_pptx(path)
    assert result.path == path.resolve()
    assert result.file_name == "sample.pptx"
    assert result.size_bytes == path.stat().st_size
    assert result.slide_count == 2
    assert len(result.digest) == 64


@pytest.mark.parametrize("name", ["sample.ppt", "sample.pptm", "sample.txt"])
def test_inspect_pptx_rejects_unsupported_extensions(tmp_path: Path, name: str) -> None:
    path = tmp_path / name
    path.write_bytes(b"not a presentation")
    with pytest.raises(PptxImportError) as caught:
        inspect_pptx(path)
    assert caught.value.code == "unsupported_format"


def test_inspect_pptx_rejects_invalid_package(tmp_path: Path) -> None:
    path = tmp_path / "bad.pptx"
    path.write_bytes(b"not a zip")
    with pytest.raises(PptxImportError) as caught:
        inspect_pptx(path)
    assert caught.value.code == "invalid_pptx"


def test_inspect_pptx_rejects_missing_required_parts(tmp_path: Path) -> None:
    path = tmp_path / "incomplete.pptx"
    with ZipFile(path, "w") as package:
        package.writestr("readme.txt", "incomplete")
    with pytest.raises(PptxImportError) as caught:
        inspect_pptx(path)
    assert caught.value.code == "invalid_pptx"
