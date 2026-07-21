from dataclasses import replace
from pathlib import Path

import pytest
from pptx import Presentation

from slideguard.pptx.importer import inspect_pptx
from slideguard.repair.planner import FixPlanError, build_fix_plan, select_fix_operations, validate_plan_source
from slideguard.rules.factory import issue
from slideguard.rules.models import IssueStatus, Severity
from slideguard.scan.models import ScanMode, ScanRequest
from slideguard.scan.orchestrator import run_scan


def _result(tmp_path: Path, *, mode: ScanMode = ScanMode.STANDARD, complete: bool = True):  # type: ignore[no-untyped-def]
    path = tmp_path / "source.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.save(path)
    kinds = (
        ("move_x", "100", "R007"),
        ("set_title_style", "style", "R009"),
        ("replace_font_size", "18", "R005"),
        ("replace_font", "Microsoft YaHei", "R004"),
    )
    issues = tuple(
        issue(
            fact_key=f"fact-{kind}",
            rule_id=rule_id,
            slide_index=1,
            object_keys=("ppt/slides/slide1.xml:shape:2",),
            severity=Severity.S3,
            actual_value="old",
            expected_value=target,
            evidence="evidence",
            suggestion="suggestion",
            can_auto_fix=True,
            fix_kind=kind,
            fix_target=target,
        )
        for kind, target, rule_id in kinds
    )
    rules = {"R002": lambda *_: issues}
    result = run_scan(
        inspect_pptx(path),
        ScanRequest(ScanMode.CUSTOM, selected_rules=("R002",)),
        rules=rules,
    )
    return replace(result, mode=mode, complete=complete)


def test_fix_plan_requires_complete_standard_scan_and_orders_operations(tmp_path: Path) -> None:
    result = _result(tmp_path)
    selected = tuple(found.issue_id for found in reversed(result.issues))
    plan = build_fix_plan(result, selected, tmp_path / "fixed.pptx")
    assert [operation.property_name for operation in plan.operations] == [
        "replace_font", "replace_font_size", "set_title_style", "move_x"
    ]
    assert plan.issue_ids == selected
    validate_plan_source(plan)


@pytest.mark.parametrize("mode,complete", [(ScanMode.QUICK, True), (ScanMode.STANDARD, False)])
def test_fix_plan_rejects_nonstandard_or_incomplete_scan(tmp_path: Path, mode: ScanMode, complete: bool) -> None:
    result = _result(tmp_path, mode=mode, complete=complete)
    with pytest.raises(FixPlanError, match="完整完成"):
        build_fix_plan(result, (result.issues[0].issue_id,), tmp_path / "fixed.pptx")


def test_fix_plan_rejects_nonfixable_issue_and_overwrite(tmp_path: Path) -> None:
    result = _result(tmp_path)
    nonfixable = replace(result.issues[0], can_auto_fix=False, fix_proposal=None)
    result = replace(result, issues=(nonfixable, *result.issues[1:]))
    with pytest.raises(FixPlanError, match="不支持"):
        build_fix_plan(result, (nonfixable.issue_id,), tmp_path / "fixed.pptx")
    destination = tmp_path / "exists.pptx"
    destination.write_bytes(b"keep")
    valid = _result(tmp_path)
    with pytest.raises(FixPlanError, match="不允许覆盖"):
        build_fix_plan(valid, (valid.issues[0].issue_id,), destination)


def test_source_hash_change_invalidates_plan(tmp_path: Path) -> None:
    result = _result(tmp_path)
    plan = build_fix_plan(result, (result.issues[0].issue_id,), tmp_path / "fixed.pptx")
    result.snapshot.file_identity.path.write_bytes(b"changed")
    with pytest.raises(FixPlanError, match="发生变化"):
        validate_plan_source(plan)


def test_fix_plan_rejects_ignored_issue(tmp_path: Path) -> None:
    result = _result(tmp_path)
    ignored = replace(result.issues[0], status=IssueStatus.IGNORED)
    result = replace(result, issues=(ignored, *result.issues[1:]))
    with pytest.raises(FixPlanError, match="不是待处理状态"):
        build_fix_plan(result, (ignored.issue_id,), tmp_path / "fixed.pptx")


def test_fix_plan_can_keep_only_confirmed_operations(tmp_path: Path) -> None:
    result = _result(tmp_path)
    plan = build_fix_plan(
        result,
        tuple(found.issue_id for found in result.issues),
        tmp_path / "fixed.pptx",
    )
    reduced = select_fix_operations(plan, (0, 2))
    assert [operation.property_name for operation in reduced.operations] == [
        "replace_font", "set_title_style"
    ]
    assert set(reduced.issue_ids) == {
        issue_id for operation in reduced.operations for issue_id in operation.issue_ids
    }
    with pytest.raises(FixPlanError, match="至少保留"):
        select_fix_operations(plan, ())
