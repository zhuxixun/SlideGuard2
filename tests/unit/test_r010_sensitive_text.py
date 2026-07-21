from pathlib import Path

from pptx import Presentation

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx.snapshot import build_snapshot
from slideguard.rules.models import IssueStatus, Severity
from slideguard.rules.r010_sensitive_text import check_sensitive_text


def test_r010_emits_stable_s1_issue_for_every_literal_match(tmp_path: Path) -> None:
    path = tmp_path / "sensitive.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000).text = "XXababaYY"
    document.save(path)
    snapshot = build_snapshot(inspect_pptx(path))

    first = check_sensitive_text(snapshot, [" aba ", "aba", "ba"])
    second = check_sensitive_text(snapshot, ["aba", "ba"])

    assert first == second
    assert [issue.actual_value for issue in first] == ["aba", "aba", "ba", "ba"]
    assert all(issue.rule_id == "R010" for issue in first)
    assert all(issue.severity is Severity.S1 for issue in first)
    assert all(issue.status is IssueStatus.PENDING for issue in first)
    assert all(issue.slide_index == 1 for issue in first)
    assert all(issue.can_auto_fix is False for issue in first)
    assert all(issue.fix_proposal is None for issue in first)
    assert len({issue.issue_id for issue in first}) == 4


def test_r010_is_case_sensitive_and_empty_lexicon_has_no_issues(tmp_path: Path) -> None:
    path = tmp_path / "case-sensitive.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000).text = "ProjectX"
    document.save(path)
    snapshot = build_snapshot(inspect_pptx(path))

    assert check_sensitive_text(snapshot, []) == ()
    assert check_sensitive_text(snapshot, ["projectx"]) == ()
    assert len(check_sensitive_text(snapshot, ["ProjectX"])) == 1
