from pathlib import Path

from fastapi.testclient import TestClient

from slideguard.lexicon import LexiconStore
from slideguard.server.app import create_app
from slideguard.server.lifecycle import LifecycleController
from slideguard.application.session import SessionStore
from slideguard.pptx.importer import inspect_pptx
from slideguard.scan.manager import ScanManager
from slideguard.scan.models import ScanMode, ScanRequest
from pptx import Presentation


TOKEN_HEADERS = {"X-SlideGuard-Token": "secret"}


def test_api_requires_session_token_and_sets_security_headers() -> None:
    client = TestClient(create_app(token="secret"))

    assert client.get("/api/health").status_code == 401
    assert client.get(
        "/api/health", headers={"X-SlideGuard-Token": "wrong"}
    ).status_code == 401
    response = client.get("/api/health", headers=TOKEN_HEADERS)
    assert response.status_code == 200
    assert response.json() == {"status": "ok"}
    assert response.headers["cache-control"] == "no-store"
    assert response.headers["x-content-type-options"] == "nosniff"
    assert "default-src 'self'" in response.headers["content-security-policy"]


def test_api_rejects_unexpected_host_and_origin() -> None:
    client = TestClient(
        create_app(
            token="secret",
            expected_host="testserver",
            allowed_origin="http://testserver",
        )
    )

    assert client.get(
        "/api/health",
        headers={**TOKEN_HEADERS, "Host": "evil.example"},
    ).status_code == 400
    assert client.get(
        "/api/health",
        headers={**TOKEN_HEADERS, "Origin": "https://evil.example"},
    ).status_code == 403
    assert client.get(
        "/api/health",
        headers={**TOKEN_HEADERS, "Origin": "http://testserver"},
    ).status_code == 200


def test_lexicon_api_reads_normalizes_and_saves(tmp_path: Path) -> None:
    path = tmp_path / "sensitive-terms.txt"
    path.write_text("旧项目\n", encoding="utf-8")
    client = TestClient(create_app(token="secret", lexicon_store=LexiconStore(path)))

    before = client.get("/api/lexicon", headers=TOKEN_HEADERS)
    assert before.status_code == 200
    assert before.json()["terms"] == ["旧项目"]

    after = client.put(
        "/api/lexicon",
        headers=TOKEN_HEADERS,
        json={
            "terms": [" 内部代号 ", "内部代号", "禁用产品"],
            "expected_digest": before.json()["digest"],
        },
    )
    assert after.status_code == 200
    assert after.json()["terms"] == ["内部代号", "禁用产品"]
    assert after.json()["count"] == 2
    assert after.json()["empty"] is False


def test_lexicon_api_rejects_stale_update(tmp_path: Path) -> None:
    path = tmp_path / "sensitive-terms.txt"
    path.write_text("词条一\n", encoding="utf-8")
    client = TestClient(create_app(token="secret", lexicon_store=LexiconStore(path)))
    digest = client.get("/api/lexicon", headers=TOKEN_HEADERS).json()["digest"]
    path.write_text("词条二\n", encoding="utf-8")

    response = client.put(
        "/api/lexicon",
        headers=TOKEN_HEADERS,
        json={"terms": ["词条三"], "expected_digest": digest},
    )
    assert response.status_code == 409
    assert response.json()["detail"]["code"] == "lexicon_conflict"


def test_websocket_requires_token_protocol_and_tracks_connection() -> None:
    lifecycle = LifecycleController(idle_seconds=60)
    client = TestClient(
        create_app(
            token="secret",
            lifecycle=lifecycle,
            allowed_origin="http://testserver",
        )
    )

    with client.websocket_connect(
        "/ws",
        subprotocols=["slideguard", "secret"],
        headers={"Origin": "http://testserver"},
    ) as websocket:
        assert websocket.accepted_subprotocol == "slideguard"
        assert lifecycle.connected_clients == 1
        websocket.send_text("ping")
        assert websocket.receive_text() == "pong"

    assert lifecycle.connected_clients == 0


def test_exit_endpoint_requests_shutdown() -> None:
    events: list[str] = []
    lifecycle = LifecycleController(idle_seconds=60)
    lifecycle.set_shutdown_callback(lambda: events.append("shutdown"))
    client = TestClient(create_app(token="secret", lifecycle=lifecycle))

    response = client.post("/api/exit", headers=TOKEN_HEADERS)

    assert response.status_code == 202
    import time

    deadline = time.monotonic() + 0.5
    while not events and time.monotonic() < deadline:
        time.sleep(0.01)
    assert events == ["shutdown"]


def test_scan_api_requires_presentation_and_returns_completed_result(tmp_path: Path) -> None:
    manager = ScanManager()
    sessions = SessionStore()
    lexicon_path = tmp_path / "terms.txt"
    lexicon_path.write_text("Secret\n", encoding="utf-8")
    client = TestClient(
        create_app(
            token="secret",
            session_store=sessions,
            scan_manager=manager,
            lexicon_store=LexiconStore(lexicon_path),
        )
    )
    missing = client.post(
        "/api/scans",
        headers=TOKEN_HEADERS,
        json={"mode": "quick"},
    )
    assert missing.status_code == 409

    pptx_path = tmp_path / "scan.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.save(pptx_path)
    sessions.replace(inspect_pptx(pptx_path))
    started = client.post(
        "/api/scans",
        headers=TOKEN_HEADERS,
        json={"mode": "quick"},
    )
    assert started.status_code == 202

    import time
    deadline = time.monotonic() + 5
    payload = {}
    while time.monotonic() < deadline:
        payload = client.get("/api/scans/current", headers=TOKEN_HEADERS).json()
        if payload["state"] != "running":
            break
        time.sleep(0.01)
    assert payload["state"] == "completed"
    assert payload["result"]["requested_rules"] == list(
        ("R002", "R003", "R004", "R006", "R009", "R010")
    )
    assert payload["result"]["rule_set_version"] == "builtin-rules-v1.0"
    issue_id = payload["result"]["issues"][0]["issue_id"]
    preview = client.get(
        f"/api/scans/current/slides/1/preview?issue_id={issue_id}",
        headers=TOKEN_HEADERS,
    )
    assert preview.status_code == 200
    assert preview.headers["content-type"].startswith("image/svg+xml")
    assert 'data-page-highlight="true"' in preview.text

    ignored = client.put(
        f"/api/scans/current/issues/{issue_id}/status",
        headers=TOKEN_HEADERS,
        json={"status": "ignored"},
    )
    assert ignored.status_code == 200
    assert ignored.json() == {"issue_id": issue_id, "status": "ignored"}
    refreshed = client.get("/api/scans/current", headers=TOKEN_HEADERS).json()
    assert refreshed["result"]["issues"][0]["status"] == "ignored"

    restored = client.put(
        f"/api/scans/current/issues/{issue_id}/status",
        headers=TOKEN_HEADERS,
        json={"status": "pending"},
    )
    assert restored.status_code == 200


def test_scan_api_rejects_empty_custom_selection(tmp_path: Path) -> None:
    path = tmp_path / "custom.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.save(path)
    sessions = SessionStore()
    sessions.replace(inspect_pptx(path))
    client = TestClient(
        create_app(token="secret", session_store=sessions, scan_manager=ScanManager())
    )
    response = client.post(
        "/api/scans",
        headers=TOKEN_HEADERS,
        json={"mode": "custom", "selected_rules": []},
    )
    assert response.status_code == 409


def test_opening_another_file_clears_stale_scan_result(tmp_path: Path) -> None:
    first = tmp_path / "first.pptx"
    second = tmp_path / "second.pptx"
    for path in (first, second):
        document = Presentation()
        document.slides.add_slide(document.slide_layouts[6])
        document.save(path)

    from slideguard.server.native_dialog import NativeDialogService

    class DialogBackend:
        def open_pptx(self) -> str:
            return str(second)

        def save_report(self, default_name: str, extension: str) -> None:
            return None

        def close(self) -> None:
            return None

    sessions = SessionStore()
    manager = ScanManager()
    dialogs = NativeDialogService(DialogBackend)
    sessions.replace(inspect_pptx(first))
    manager.start(sessions.current().presentation, ScanRequest(ScanMode.QUICK))
    import time
    deadline = time.monotonic() + 5
    while manager.snapshot().state.value == "running" and time.monotonic() < deadline:
        time.sleep(0.01)
    assert manager.snapshot().result is not None

    with TestClient(
        create_app(
            token="secret",
            session_store=sessions,
            scan_manager=manager,
            native_dialog=dialogs,
        )
    ) as client:
        response = client.post("/api/dialog/open-pptx", headers=TOKEN_HEADERS)
        assert response.status_code == 200
        assert client.get("/api/scans/current", headers=TOKEN_HEADERS).json()["state"] == "idle"
        assert sessions.current().presentation.path == second.resolve()
    assert sessions.current() is None


def test_dropped_file_is_imported_locally_and_cleaned_with_session(tmp_path: Path) -> None:
    source = tmp_path / "拖入 样例.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.save(source)

    class DialogStub:
        def close(self) -> None:
            return None

    sessions = SessionStore()
    managed_path = None
    with TestClient(
        create_app(
            token="secret",
            session_store=sessions,
            scan_manager=ScanManager(),
            native_dialog=DialogStub(),  # type: ignore[arg-type]
            import_dir=tmp_path / "sessions",
        )
    ) as client:
        response = client.post(
            "/api/files/drop",
            headers={**TOKEN_HEADERS, "X-SlideGuard-Filename": "%E6%8B%96%E5%85%A5%20%E6%A0%B7%E4%BE%8B.pptx"},
            content=source.read_bytes(),
        )
        assert response.status_code == 200
        assert response.json()["file"]["name"] == "拖入 样例.pptx"
        managed_path = Path(response.json()["file"]["path"])
        assert managed_path.is_file()
        assert sessions.current().managed_copy is True
    assert managed_path is not None
    assert not managed_path.exists()


def test_scan_continues_when_sensitive_lexicon_is_invalid(tmp_path: Path) -> None:
    pptx_path = tmp_path / "invalid-lexicon.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.save(pptx_path)
    sessions = SessionStore()
    sessions.replace(inspect_pptx(pptx_path))
    lexicon_path = tmp_path / "terms.txt"
    lexicon_path.write_bytes(b"\xff\xfe")
    manager = ScanManager()
    client = TestClient(
        create_app(
            token="secret",
            session_store=sessions,
            scan_manager=manager,
            lexicon_store=LexiconStore(lexicon_path),
        )
    )
    response = client.post("/api/scans", headers=TOKEN_HEADERS, json={"mode": "standard"})
    assert response.status_code == 202
    import time
    deadline = time.monotonic() + 5
    payload = {}
    while time.monotonic() < deadline:
        payload = client.get("/api/scans/current", headers=TOKEN_HEADERS).json()
        if payload["state"] != "running":
            break
        time.sleep(0.01)
    assert payload["state"] == "incomplete"
    assert payload["result"]["failures"][0]["rule_id"] == "R010"
    assert "R002" in payload["result"]["completed_rules"]
