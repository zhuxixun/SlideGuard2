from dataclasses import replace
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx.snapshot import build_snapshot
from slideguard.rules.r009_title import check_titles


def _set_title(slide, text: str):  # type: ignore[no-untyped-def]
    title = slide.shapes.title
    title.text_frame.clear()
    run = title.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(192, 0, 0)
    return title


def _save_snapshot(document: Presentation, path: Path):  # type: ignore[no-untyped-def]
    document.save(path)
    return build_snapshot(inspect_pptx(path))


def test_r009_accepts_compliant_title_and_ignores_normal_textbox(tmp_path: Path) -> None:
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[0])
    _set_title(slide, "合规标题")
    box = slide.shapes.add_textbox(0, 0, Inches(4), Inches(1))
    box.text = "普通文本框即使样式不同也不是标题"
    snapshot = _save_snapshot(document, tmp_path / "compliant.pptx")
    assert check_titles(snapshot) == ()


def test_r009_reports_style_as_one_fixable_issue(tmp_path: Path) -> None:
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Wrong title"
    snapshot = _save_snapshot(document, tmp_path / "wrong-style.pptx")
    issues = check_titles(snapshot)
    assert len(issues) == 1
    assert "字体" in issues[0].actual_value
    assert "字号" in issues[0].actual_value
    assert "加粗" in issues[0].actual_value
    assert "颜色" in issues[0].actual_value
    assert issues[0].fix_proposal.kind == "set_title_style"


def test_r009_prefers_layout_title_reference_line(tmp_path: Path) -> None:
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[0])
    title = _set_title(slide, "位置检查")
    title.left = title.left + Pt(4)
    snapshot = _save_snapshot(document, tmp_path / "position.pptx")
    issues = check_titles(snapshot)
    assert len(issues) == 1
    assert "偏差 +4.00pt" in issues[0].evidence
    assert issues[0].fix_proposal.kind == "move_x"


def test_r009_scales_only_suffix_after_first_colon(tmp_path: Path) -> None:
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[0])
    title = _set_title(slide, "主题：这是一段需要缩小字号才能放入单行的标题文字")
    title.width = Inches(5.5)
    snapshot = _save_snapshot(document, tmp_path / "colon.pptx")
    issues = check_titles(snapshot)
    overflow = next(item for item in issues if item.fact_key.endswith("title-overflow"))
    assert overflow.can_auto_fix is True
    assert overflow.fix_proposal.kind == "scale_title_suffix"
    assert 14 <= int(overflow.fix_proposal.target_value) <= 23


def test_r009_long_title_without_colon_requires_manual_handling(tmp_path: Path) -> None:
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[0])
    title = _set_title(slide, "这是一个没有冒号而且长度明显超过标题区域无法单行展示的标题")
    title.width = Inches(3)
    snapshot = _save_snapshot(document, tmp_path / "no-colon.pptx")
    issues = check_titles(snapshot)
    overflow = next(item for item in issues if item.fact_key.endswith("title-overflow"))
    assert overflow.can_auto_fix is False
    assert overflow.fix_proposal is None


def test_r009_falls_back_to_seventy_percent_reference(tmp_path: Path) -> None:
    document = Presentation()
    for index in range(4):
        slide = document.slides.add_slide(document.slide_layouts[0])
        _set_title(slide, f"标题 {index + 1}")
    snapshot = _save_snapshot(document, tmp_path / "fallback.pptx")
    slides = []
    for index, slide in enumerate(snapshot.slides):
        obj = slide.objects[0]
        if index == 3:
            obj = replace(obj, bounds_pt=replace(obj.bounds_pt, left=obj.bounds_pt.left + 10))
        slides.append(replace(slide, layout_title_left_pt=None, objects=(obj, *slide.objects[1:])))
    issues = check_titles(replace(snapshot, slides=tuple(slides)))
    assert len(issues) == 1
    assert issues[0].slide_index == 4
    assert issues[0].fix_proposal.kind == "move_x"
