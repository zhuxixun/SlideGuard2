from dataclasses import replace
from pathlib import Path

import pytest
from openpyxl import load_workbook
from pptx import Presentation

from slideguard.pptx.importer import inspect_pptx
from slideguard.reporting.exporters import default_report_name, export_html, export_xlsx
from slideguard.rules.factory import issue
from slideguard.rules.models import Severity
from slideguard.scan.models import RepairComparison, ScanMode, ScanRequest
from slideguard.scan.orchestrator import run_scan


def _result(tmp_path: Path, *, complete: bool = True):  # type: ignore[no-untyped-def]
    path = tmp_path / "客户 & 项目.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.save(path)
    found = issue(
        fact_key="report-test",
        rule_id="R002",
        slide_index=1,
        object_keys=(),
        severity=Severity.S1,
        actual_value="=SUM(A1:A2)<script>",
        expected_value="expected",
        evidence="evidence",
        suggestion="suggestion",
    )
    unavailable = {"R003": "failed"} if not complete else None
    rules = {"R002": lambda *_: (found,), "R003": lambda *_: ()}
    selected = ("R002", "R003") if not complete else ("R002",)
    return run_scan(
        inspect_pptx(path),
        ScanRequest(ScanMode.CUSTOM, selected_rules=selected),
        rules=rules,
        unavailable_rules=unavailable,
    )


def test_html_report_is_single_file_escaped_and_marks_incomplete(tmp_path: Path) -> None:
    result = _result(tmp_path, complete=False)
    destination = tmp_path / "report.html"
    export_html(result, destination)
    content = destination.read_text(encoding="utf-8")
    assert "扫描未完成" in content
    assert "请求执行规则" in content
    assert "R003：failed" in content
    assert "客户 &amp; 项目.pptx" in content
    assert "&lt;script&gt;" in content
    assert "=SUM(A1:A2)&lt;script&gt;" in content
    assert "report-search" in content
    assert "report-severity" in content
    assert "report-rule" in content
    assert "addEventListener('input', apply)" in content
    assert "data-page-highlight" in content
    assert 'src="http' not in content
    assert 'href="http' not in content


def test_xlsx_report_has_summary_details_and_formula_protection(tmp_path: Path) -> None:
    destination = tmp_path / "report.xlsx"
    export_xlsx(_result(tmp_path), destination)
    workbook = load_workbook(destination, data_only=False)
    assert workbook.sheetnames == ["扫描摘要", "问题清单"]
    details = workbook["问题清单"]
    assert details.cell(2, 7).value == "'=SUM(A1:A2)<script>"
    assert workbook["扫描摘要"].cell(2, 2).value == "客户 & 项目.pptx"
    summary = {
        row[0].value: row[1].value
        for row in workbook["扫描摘要"].iter_rows(min_row=2)
    }
    assert summary["请求执行规则"] == "R002"
    assert summary["涉及问题页面数"] == 1
    assert summary["可自动修复问题数"] == 0
    assert summary["R002 问题数"] == 1
    assert details.cell(2, 13).value == "否"


def test_report_export_never_overwrites_existing_file(tmp_path: Path) -> None:
    destination = tmp_path / "existing.html"
    destination.write_text("keep", encoding="utf-8")
    with pytest.raises(FileExistsError):
        export_html(_result(tmp_path), destination)
    assert destination.read_text(encoding="utf-8") == "keep"


def test_default_report_name_contains_source_and_timestamp(tmp_path: Path) -> None:
    name = default_report_name(_result(tmp_path), "xlsx")
    assert name.startswith("客户 & 项目_SlideGuard_")
    assert name.endswith(".xlsx")


def test_reports_include_repair_comparison(tmp_path: Path) -> None:
    result = replace(
        _result(tmp_path),
        repair_comparison=RepairComparison(3, 2, 1, 1),
    )
    html_path = tmp_path / "comparison.html"
    xlsx_path = tmp_path / "comparison.xlsx"
    export_html(result, html_path)
    export_xlsx(result, xlsx_path)

    content = html_path.read_text(encoding="utf-8")
    assert "修复前后对比" in content
    assert "已修复 2" in content
    summary = load_workbook(xlsx_path)["扫描摘要"]
    values = {row[0].value: row[1].value for row in summary.iter_rows(min_row=2)}
    assert values["修复选中问题数"] == 3
    assert values["修复后新增问题数"] == 1


def test_reports_warn_when_sensitive_lexicon_is_empty(tmp_path: Path) -> None:
    result = replace(_result(tmp_path), sensitive_lexicon_empty=True)
    html_path = tmp_path / "empty-lexicon.html"
    xlsx_path = tmp_path / "empty-lexicon.xlsx"
    export_html(result, html_path)
    export_xlsx(result, xlsx_path)
    assert "敏感词库为空，本项未发现问题不代表无敏感内容" in html_path.read_text(encoding="utf-8")
    summary = load_workbook(xlsx_path)["扫描摘要"]
    values = {row[0].value: row[1].value for row in summary.iter_rows(min_row=2)}
    assert values["敏感词库状态"].startswith("为空")
