from dataclasses import FrozenInstanceError
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx import snapshot as snapshot_module
from slideguard.pptx.snapshot import ParseStatus, TextSource, build_snapshot


def test_build_snapshot_captures_geometry_text_and_character_locations(tmp_path: Path) -> None:
    path = tmp_path / "snapshot.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Internal Name"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(24)
    run.font.bold = True
    document.save(path)

    snapshot = build_snapshot(inspect_pptx(path))

    assert snapshot.slide_width_pt == pytest.approx(720)
    assert snapshot.slide_height_pt == pytest.approx(540)
    assert len(snapshot.slides) == 1
    shape = snapshot.slides[0].objects[0]
    assert shape.bounds_pt.left == pytest.approx(72)
    assert shape.bounds_pt.top == pytest.approx(144)
    assert shape.text_frame.text == "Internal Name"
    assert shape.text_frame.paragraphs[0][0].font_name == "Microsoft YaHei"
    assert shape.text_frame.paragraphs[0][0].font_size_pt == pytest.approx(24)
    occurrence = next(item for item in snapshot.text_occurrences if item.text == "Internal Name")
    assert occurrence.source is TextSource.SLIDE
    assert occurrence.slide_index == 1
    assert len(occurrence.character_map) == len(occurrence.text)
    assert occurrence.character_map[3].offset == 3


def test_snapshot_is_immutable(tmp_path: Path) -> None:
    path = tmp_path / "immutable.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.save(path)
    snapshot = build_snapshot(inspect_pptx(path))
    with pytest.raises(FrozenInstanceError):
        snapshot.slide_width_pt = 1  # type: ignore[misc]


def test_snapshot_preserves_hidden_slide_state(tmp_path: Path) -> None:
    path = tmp_path / "hidden.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.slides._sldIdLst[0].set("show", "0")  # noqa: SLF001
    document.save(path)
    snapshot = build_snapshot(inspect_pptx(path))
    assert snapshot.slides[0].hidden is True


def test_snapshot_skips_one_failed_object_and_marks_page_partial(tmp_path: Path, monkeypatch) -> None:  # type: ignore[no-untyped-def]
    path = tmp_path / "partial.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text = "first"
    slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1)).text = "second"
    document.save(path)
    original = snapshot_module._shape_snapshot
    calls = 0

    def fail_first(*args, **kwargs):  # type: ignore[no-untyped-def]
        nonlocal calls
        calls += 1
        if calls == 1:
            raise ValueError("must not be exposed")
        return original(*args, **kwargs)

    monkeypatch.setattr(snapshot_module, "_shape_snapshot", fail_first)

    snapshot = build_snapshot(inspect_pptx(path))

    assert len(snapshot.slides[0].objects) == 1
    assert snapshot.slides[0].parse_status is ParseStatus.PARTIAL
    assert snapshot.parse_failures[0].scope == "object"
    assert snapshot.parse_failures[0].reason == "ValueError"
