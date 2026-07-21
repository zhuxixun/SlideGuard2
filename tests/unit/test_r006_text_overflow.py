from dataclasses import replace
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx.snapshot import build_snapshot
from slideguard.preview.text_flow import TextFlowMeasurement
from slideguard.rules.models import Severity
from slideguard.rules.r006_text_overflow import check_text_overflow


def _snapshot(tmp_path: Path, text: str, font: str = "Microsoft YaHei"):  # type: ignore[no-untyped-def]
    path = tmp_path / "overflow.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, Inches(2), Inches(1))
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.word_wrap = False
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(20)
    document.save(path)
    return build_snapshot(inspect_pptx(path))


def test_r006_reports_measured_width_overflow(tmp_path: Path) -> None:
    snapshot = _snapshot(tmp_path, "这是一段明显超过文本框宽度的微软雅黑文字")
    issues = check_text_overflow(snapshot)
    assert len(issues) == 1
    assert issues[0].severity is Severity.S3
    assert "宽度超出" in issues[0].evidence
    assert issues[0].can_auto_fix is False


def test_r006_applies_two_point_tolerance(tmp_path: Path, monkeypatch) -> None:  # type: ignore[no-untyped-def]
    snapshot = _snapshot(tmp_path, "Tolerance")
    obj = snapshot.slides[0].objects[0]
    frame = replace(obj.text_frame, margin_left_pt=0, margin_right_pt=0)
    obj = replace(obj, bounds_pt=replace(obj.bounds_pt, width=100, height=100), text_frame=frame)
    snapshot = replace(snapshot, slides=(replace(snapshot.slides[0], objects=(obj,)),))

    monkeypatch.setattr(
        "slideguard.rules.r006_text_overflow.measure_text_flow",
        lambda *_: TextFlowMeasurement(102, 20, True, ()),
    )
    assert check_text_overflow(snapshot) == ()

    monkeypatch.setattr(
        "slideguard.rules.r006_text_overflow.measure_text_flow",
        lambda *_: TextFlowMeasurement(102.01, 20, True, ()),
    )
    assert len(check_text_overflow(snapshot)) == 1


def test_r006_does_not_claim_overflow_for_unavailable_font(tmp_path: Path) -> None:
    snapshot = _snapshot(tmp_path, "A very long Arial sentence", font="Arial")
    assert check_text_overflow(snapshot) == ()


def test_r006_marks_text_beyond_page_as_s2(tmp_path: Path) -> None:
    snapshot = _snapshot(tmp_path, "页面外文字")
    obj = snapshot.slides[0].objects[0]
    obj = replace(obj, bounds_pt=replace(obj.bounds_pt, left=snapshot.slide_width_pt - 5))
    snapshot = replace(snapshot, slides=(replace(snapshot.slides[0], objects=(obj,)),))
    issues = check_text_overflow(snapshot)
    assert len(issues) == 1
    assert issues[0].severity is Severity.S2
    assert "文字超出页面" in issues[0].evidence
