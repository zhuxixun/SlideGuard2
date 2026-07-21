from pathlib import Path

import pytest
from pptx import Presentation

from slideguard.pptx.importer import inspect_pptx
from slideguard.rules.factory import issue
from slideguard.rules.models import Severity
from slideguard.scan.models import ScanMode, ScanRequest, ScanStage
from slideguard.scan.orchestrator import (
    ALL_RULES,
    QUICK_RULES,
    CancellationToken,
    run_scan,
)


def _imported(tmp_path: Path):  # type: ignore[no-untyped-def]
    path = tmp_path / "scan.pptx"
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.save(path)
    return inspect_pptx(path)


def _recording_rules(rule_ids, calls):  # type: ignore[no-untyped-def]
    def make(rule_id):  # type: ignore[no-untyped-def]
        def check(snapshot, terms):  # type: ignore[no-untyped-def]
            calls.append((rule_id, terms))
            return ()
        return check
    return {rule_id: make(rule_id) for rule_id in rule_ids}


def test_quick_scan_runs_exact_required_rules_and_progress_stages(tmp_path: Path) -> None:
    calls: list[tuple[str, tuple[str, ...]]] = []
    events = []
    result = run_scan(
        _imported(tmp_path),
        ScanRequest(ScanMode.QUICK, sensitive_terms=("Secret",)),
        rules=_recording_rules(QUICK_RULES, calls),
        on_progress=events.append,
    )
    assert result.requested_rules == QUICK_RULES
    assert result.completed_rules == QUICK_RULES
    assert result.complete is True
    assert [rule_id for rule_id, _ in calls] == list(QUICK_RULES)
    assert all(terms == ("Secret",) for _, terms in calls)
    assert [event.stage for event in events if event.stage is not ScanStage.CHECKING] == [
        ScanStage.PARSING,
        ScanStage.PREVIEW,
        ScanStage.SUMMARIZING,
    ]


def test_standard_scan_runs_all_base_rules(tmp_path: Path) -> None:
    calls = []
    result = run_scan(
        _imported(tmp_path),
        ScanRequest(ScanMode.STANDARD),
        rules=_recording_rules(ALL_RULES, calls),
    )
    assert result.requested_rules == ALL_RULES
    assert [rule_id for rule_id, _ in calls] == list(ALL_RULES)


def test_custom_scan_orders_and_deduplicates_selection(tmp_path: Path) -> None:
    calls = []
    result = run_scan(
        _imported(tmp_path),
        ScanRequest(ScanMode.CUSTOM, selected_rules=("R010", "R002", "R010")),
        rules=_recording_rules(("R002", "R010"), calls),
    )
    assert result.requested_rules == ("R002", "R010")
    with pytest.raises(ValueError, match="至少选择"):
        run_scan(_imported(tmp_path), ScanRequest(ScanMode.CUSTOM))


def test_rule_failure_is_isolated_and_marks_result_incomplete(tmp_path: Path) -> None:
    calls = []
    rules = _recording_rules(("R002", "R003"), calls)
    rules["R002"] = lambda *_: (_ for _ in ()).throw(RuntimeError("broken rule"))
    result = run_scan(
        _imported(tmp_path),
        ScanRequest(ScanMode.CUSTOM, selected_rules=("R002", "R003")),
        rules=rules,
    )
    assert result.completed_rules == ("R003",)
    assert result.failures[0].rule_id == "R002"
    assert result.failures[0].message == "broken rule"
    assert result.complete is False
    assert calls == [("R003", ())]


def test_preflight_rule_failure_does_not_block_other_rules(tmp_path: Path) -> None:
    calls = []
    result = run_scan(
        _imported(tmp_path),
        ScanRequest(ScanMode.CUSTOM, selected_rules=("R002", "R010")),
        rules=_recording_rules(("R002", "R010"), calls),
        unavailable_rules={"R010": "词库不是有效 UTF-8"},
    )
    assert calls == [("R002", ())]
    assert result.completed_rules == ("R002",)
    assert result.failures[0].rule_id == "R010"
    assert result.complete is False


def test_cancellation_stops_scheduling_new_rules(tmp_path: Path) -> None:
    token = CancellationToken()
    calls = []

    def first(*_):  # type: ignore[no-untyped-def]
        calls.append("R002")
        token.cancel()
        return ()

    result = run_scan(
        _imported(tmp_path),
        ScanRequest(ScanMode.CUSTOM, selected_rules=("R002", "R003")),
        cancellation=token,
        rules={"R002": first, "R003": lambda *_: calls.append("R003") or ()},
    )
    assert calls == ["R002"]
    assert result.completed_rules == ("R002",)
    assert result.cancelled is True
    assert result.complete is False


def test_issues_are_deduplicated_and_sorted(tmp_path: Path) -> None:
    def found(rule_id: str, severity: Severity, slide: int, fact: str):
        return issue(
            fact_key=fact,
            rule_id=rule_id,
            slide_index=slide,
            object_keys=(),
            severity=severity,
            actual_value="actual",
            expected_value="expected",
            evidence="evidence",
            suggestion="suggestion",
        )

    duplicate = found("R003", Severity.S3, 1, "same")
    rules = {
        "R002": lambda *_: (found("R002", Severity.S1, 2, "critical"), duplicate),
        "R003": lambda *_: (duplicate,),
    }
    result = run_scan(
        _imported(tmp_path),
        ScanRequest(ScanMode.CUSTOM, selected_rules=("R002", "R003")),
        rules=rules,
    )
    assert [item.fact_key for item in result.issues] == ["critical", "same"]
