import asyncio
from pathlib import Path
from threading import get_ident

from fastapi.testclient import TestClient

from slideguard.server.app import create_app
from slideguard.server.native_dialog import NativeDialogService
from slideguard.application.session import SessionStore
from slideguard.scan.manager import ScanManager
from slideguard.repair.manager import RepairManager
from slideguard.pptx.importer import inspect_pptx
from pptx import Presentation


class FakeDialogBackend:
    def __init__(self, values: list[str | None]) -> None:
        self.values = values
        self.thread_ids: list[int] = []
        self.closed = False

    def open_pptx(self) -> str | None:
        self.thread_ids.append(get_ident())
        return self.values.pop(0)

    def save_report(self, default_name: str, extension: str) -> str | None:
        self.thread_ids.append(get_ident())
        return self.values.pop(0)

    def close(self) -> None:
        self.thread_ids.append(get_ident())
        self.closed = True


def test_dialog_commands_run_serially_on_owned_thread() -> None:
    backend = FakeDialogBackend(["C:/samples/one.pptx", None])
    service = NativeDialogService(lambda: backend)

    async def scenario() -> tuple[Path | None, Path | None]:
        first, second = await asyncio.gather(service.open_pptx(), service.open_pptx())
        return first, second

    try:
        assert asyncio.run(scenario()) == (Path("C:/samples/one.pptx"), None)
    finally:
        service.close()

    assert backend.closed is True
    assert len(set(backend.thread_ids)) == 1
    assert backend.thread_ids[0] != get_ident()


def test_save_report_uses_same_owned_dialog_thread(tmp_path: Path) -> None:
    destination = tmp_path / "report.html"
    backend = FakeDialogBackend([str(destination)])
    service = NativeDialogService(lambda: backend)
    try:
        selected = asyncio.run(service.save_report("default.html", ".html"))
    finally:
        service.close()
    assert selected == destination
    assert backend.closed is True
    assert len(set(backend.thread_ids)) == 1


def test_dialog_api_imports_selected_file_and_cancel(tmp_path: Path) -> None:
    sample = tmp_path / "one.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(sample)
    backend = FakeDialogBackend([str(sample), None])
    service = NativeDialogService(lambda: backend)
    sessions = SessionStore()
    headers = {"X-SlideGuard-Token": "secret"}
    with TestClient(
        create_app(token="secret", native_dialog=service, session_store=sessions)
    ) as client:
        selected = client.post("/api/dialog/open-pptx", headers=headers)
        cancelled = client.post("/api/dialog/open-pptx", headers=headers)

    assert selected.json()["cancelled"] is False
    assert selected.json()["file"]["name"] == "one.pptx"
    assert selected.json()["file"]["slide_count"] == 1
    assert cancelled.json() == {"cancelled": True, "file": None}
    assert sessions.current().presentation.path == sample.resolve()
    assert backend.closed is True


def test_repair_api_prepares_executes_and_preserves_source(tmp_path: Path) -> None:
    source = tmp_path / "repair-source.pptx"
    output = tmp_path / "repair-output.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(914400, 914400, 3657600, 914400)
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Arial text"
    run.font.name = "Arial"
    presentation.save(source)
    source_before = source.read_bytes()
    sessions = SessionStore()
    sessions.replace(inspect_pptx(source))
    scans = ScanManager()
    repairs = RepairManager()
    backend = FakeDialogBackend([str(output)])
    dialogs = NativeDialogService(lambda: backend)
    headers = {"X-SlideGuard-Token": "secret"}
    with TestClient(
        create_app(
            token="secret",
            session_store=sessions,
            scan_manager=scans,
            repair_manager=repairs,
            native_dialog=dialogs,
        )
    ) as client:
        assert client.post("/api/scans", headers=headers, json={"mode": "standard"}).status_code == 202
        import time
        deadline = time.monotonic() + 5
        state = {}
        while time.monotonic() < deadline:
            state = client.get("/api/scans/current", headers=headers).json()
            if state["state"] != "running":
                break
            time.sleep(0.01)
        font_issue = next(item for item in state["result"]["issues"] if item["rule_id"] == "R004")
        prepared = client.post(
            "/api/repairs/prepare",
            headers=headers,
            json={"issue_ids": [font_issue["issue_id"]]},
        )
        assert prepared.status_code == 200
        assert prepared.json()["plan"]["issue_count"] == 1
        executed = client.post("/api/repairs/execute", headers=headers)
        assert executed.status_code == 200
        assert executed.json()["fixed_count"] == 1
        assert executed.json()["scan"]["result"]["mode"] == "standard"
        assert executed.json()["scan"]["result"]["repair_comparison"] == {
            "selected_count": 1,
            "fixed_count": 1,
            "unresolved_count": 0,
            "introduced_count": executed.json()["introduced_count"],
        }
        refreshed = client.get("/api/scans/current", headers=headers).json()
        assert refreshed["state"] == "completed"
        assert all(
            item["issue_id"] != font_issue["issue_id"]
            for item in refreshed["result"]["issues"]
        )
        assert refreshed["result"]["requested_rules"] == [
            f"R{number:03d}" for number in range(2, 11)
        ]
    assert output.is_file()
    assert source.read_bytes() == source_before
    repaired = Presentation(output)
    assert repaired.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.name == "Microsoft YaHei"
