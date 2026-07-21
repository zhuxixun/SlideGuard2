from pathlib import Path

from pptx import Presentation

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx.snapshot import build_snapshot
from slideguard.rules.models import Severity
from slideguard.rules.r004_font import check_fonts


def _snapshot_with_runs(tmp_path: Path, runs: list[tuple[str, str]]):  # type: ignore[no-untyped-def]
    path = tmp_path / "fonts.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(0, 0, 3_000_000, 1_000_000).text_frame.paragraphs[0]
    for text, font in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = font
    document.save(path)
    return build_snapshot(inspect_pptx(path))


def test_r004_accepts_only_two_yahei_names_and_merges_contiguous_range(tmp_path: Path) -> None:
    snapshot = _snapshot_with_runs(
        tmp_path,
        [
            ("中文", "微软雅黑"),
            ("Good", "Microsoft YaHei"),
            ("Bad", "Arial"),
            ("UI", "Microsoft YaHei UI"),
        ],
    )
    issues = check_fonts(snapshot)
    assert len(issues) == 2
    assert [item.severity for item in issues] == [Severity.S3, Severity.S3]
    assert "Arial" in issues[0].actual_value
    assert "Microsoft YaHei UI" in issues[1].actual_value
    assert all(item.can_auto_fix for item in issues)
    assert all(item.fix_proposal.target_value == "Microsoft YaHei" for item in issues)


def test_r004_empty_text_box_has_no_issue(tmp_path: Path) -> None:
    snapshot = _snapshot_with_runs(tmp_path, [])
    assert check_fonts(snapshot) == ()


def test_r004_uses_theme_font_when_run_has_no_direct_font(tmp_path: Path) -> None:
    path = tmp_path / "theme-font.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 3_000_000, 1_000_000).text = "Theme text"
    document.save(path)
    issues = check_fonts(build_snapshot(inspect_pptx(path)))
    assert len(issues) == 1
    assert "Aptos" in issues[0].actual_value or "Calibri" in issues[0].actual_value
