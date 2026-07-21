from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx.snapshot import build_snapshot
from slideguard.rules.r005_font_size import check_font_sizes


def _add_text(slide, text: str, size: float):  # type: ignore[no-untyped-def]
    box = slide.shapes.add_textbox(0, 0, Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    return box


def test_r005_checks_body_minimum_but_leaves_title_to_r009(tmp_path: Path) -> None:
    path = tmp_path / "minimum.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[0])
    title_run = slide.shapes.title.text_frame.paragraphs[0].add_run()
    title_run.text = "Tiny title"
    title_run.font.size = Pt(8)
    _add_text(slide, "Tiny body", 12)
    document.save(path)

    issues = check_font_sizes(build_snapshot(inspect_pptx(path)))

    assert len(issues) == 1
    assert issues[0].actual_value == "12pt"
    assert issues[0].expected_value == "不小于 14pt"


def test_r005_uses_ten_point_minimum_for_table_text(tmp_path: Path) -> None:
    path = tmp_path / "table.pptx"
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    table = slide.shapes.add_table(1, 1, 0, 0, Inches(3), Inches(1)).table
    run = table.cell(0, 0).text_frame.paragraphs[0].add_run()
    run.text = "Small"
    run.font.size = Pt(9)
    document.save(path)
    issues = check_font_sizes(build_snapshot(inspect_pptx(path)))
    assert len(issues) == 1
    assert issues[0].expected_value == "不小于 10pt"


def test_r005_reports_mainstream_deviation_only_with_enough_support(tmp_path: Path) -> None:
    path = tmp_path / "mainstream.pptx"
    document = Presentation()
    for index, size in enumerate((18, 18, 18, 24), start=1):
        slide = document.slides.add_slide(document.slide_layouts[6])
        _add_text(slide, f"Body {index}", size)
    document.save(path)

    issues = check_font_sizes(build_snapshot(inspect_pptx(path)))

    assert len(issues) == 1
    assert issues[0].slide_index == 4
    assert issues[0].actual_value == "24pt"
    assert "主流字号 18pt" in issues[0].expected_value
    assert issues[0].fix_proposal.target_value == "18"


def test_r005_does_not_infer_mainstream_below_three_samples(tmp_path: Path) -> None:
    path = tmp_path / "two-samples.pptx"
    document = Presentation()
    for size in (18, 24):
        slide = document.slides.add_slide(document.slide_layouts[6])
        _add_text(slide, "Body", size)
    document.save(path)
    assert check_font_sizes(build_snapshot(inspect_pptx(path))) == ()
