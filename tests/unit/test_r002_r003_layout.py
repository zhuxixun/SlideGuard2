from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx.snapshot import build_snapshot
from slideguard.rules.models import Severity
from slideguard.rules.r002_blank_slide import check_blank_slides
from slideguard.rules.r003_off_slide import check_off_slide_objects


def _snapshot(document: Presentation, path: Path):  # type: ignore[no-untyped-def]
    document.save(path)
    return build_snapshot(inspect_pptx(path))


def test_r002_reports_empty_slide_but_excludes_hidden_slide(tmp_path: Path) -> None:
    document = Presentation()
    document.slides.add_slide(document.slide_layouts[6])
    document.slides.add_slide(document.slide_layouts[6])
    document.slides._sldIdLst[1].set("show", "0")  # noqa: SLF001
    issues = check_blank_slides(_snapshot(document, tmp_path / "blank.pptx"))
    assert len(issues) == 1
    assert issues[0].slide_index == 1
    assert issues[0].severity is Severity.S3
    assert issues[0].can_auto_fix is False


def test_r002_ignores_logo_and_off_canvas_but_counts_visible_shape(tmp_path: Path) -> None:
    document = Presentation()
    logo_slide = document.slides.add_slide(document.slide_layouts[6])
    logo = logo_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(1), Inches(1))
    logo.name = "Company Logo"
    off_slide = document.slides.add_slide(document.slide_layouts[6])
    off_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(20), Inches(20), Inches(1), Inches(1)
    )
    content_slide = document.slides.add_slide(document.slide_layouts[6])
    content_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(1), Inches(1))
    issues = check_blank_slides(_snapshot(document, tmp_path / "subjects.pptx"))
    assert [item.slide_index for item in issues] == [1, 2]


def test_r003_reports_only_fully_off_canvas_top_level_objects(tmp_path: Path) -> None:
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(20), Inches(20), Inches(2), Inches(1))
    text.text = "Residual text"
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(-3), Inches(1), Inches(1), Inches(1)
    )
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(-0.5), Inches(2), Inches(1), Inches(1)
    )
    issues = check_off_slide_objects(_snapshot(document, tmp_path / "outside.pptx"))
    assert len(issues) == 2
    assert [item.severity for item in issues] == [Severity.S2, Severity.S3]
    assert all(item.can_auto_fix is False for item in issues)


def test_r003_uses_rotated_visible_bounds(tmp_path: Path) -> None:
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(-1.1), Inches(1), Inches(1), Inches(3)
    )
    shape.rotation = 45
    issues = check_off_slide_objects(_snapshot(document, tmp_path / "rotation.pptx"))
    assert issues == ()
