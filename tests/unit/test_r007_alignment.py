from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from slideguard.pptx.importer import inspect_pptx
from slideguard.pptx.snapshot import build_snapshot
from slideguard.rules.r007_alignment import check_alignment


def _document_with_rectangles(positions: list[tuple[float, float]]) -> Presentation:
    document = Presentation()
    slide = document.slides.add_slide(document.slide_layouts[6])
    for left, top in positions:
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(1),
            Inches(0.5),
        )
    return document


def _check(document: Presentation, path: Path):  # type: ignore[no-untyped-def]
    document.save(path)
    return check_alignment(build_snapshot(inspect_pptx(path)))


def test_r007_selects_one_reference_and_reports_outlier(tmp_path: Path) -> None:
    issues = _check(
        _document_with_rectangles([(1, 2), (3, 2), (5, 2), (7, 2.15)]),
        tmp_path / "horizontal.pptx",
    )
    assert len(issues) == 1
    assert issues[0].slide_index == 1
    assert "top=144.00pt" in issues[0].expected_value
    assert issues[0].fix_proposal.kind == "move_y"
    assert issues[0].fix_proposal.target_value == "144.00"


def test_r007_allows_three_point_visual_adjustment(tmp_path: Path) -> None:
    three_points = 3 / 72
    issues = _check(
        _document_with_rectangles([(1, 2), (3, 2), (5, 2), (7, 2 + three_points)]),
        tmp_path / "tolerance.pptx",
    )
    assert issues == ()


def test_r007_requires_seventy_percent_support(tmp_path: Path) -> None:
    issues = _check(
        _document_with_rectangles([(1, 1), (3, 1), (5, 2)]),
        tmp_path / "support.pptx",
    )
    assert issues == ()


def test_r007_excludes_diagonal_layout(tmp_path: Path) -> None:
    issues = _check(
        _document_with_rectangles([(1, 1), (2, 2), (3, 3), (4, 4)]),
        tmp_path / "diagonal.pptx",
    )
    assert issues == ()


def test_r007_excludes_objects_with_size_difference_over_ten_percent(tmp_path: Path) -> None:
    document = _document_with_rectangles([(1, 2), (3, 2), (5, 2)])
    slide = document.slides[0]
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2.2), Inches(2), Inches(1))
    assert _check(document, tmp_path / "different-size.pptx") == ()
